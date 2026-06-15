"""Build T2_slides.pptx with proper OMML (PowerPoint native) equations.
White background, black text, minimal styling."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import os

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_M   = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
NS_A14 = 'http://schemas.microsoft.com/office/drawing/2010/main'

BLACK = RGBColor(0, 0, 0)
GREY  = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ═════════════════════════════════════════════════════════════════════════════
# OMML helpers  (all return XML strings; namespace declared at embed time)
# ═════════════════════════════════════════════════════════════════════════════

def r(t):
    """Italic math run."""
    return f'<m:r><m:t xml:space="preserve">{t}</m:t></m:r>'

def rp(t):
    """Upright (plain) run."""
    return f'<m:r><m:rPr><m:sty m:val="p"/></m:rPr><m:t xml:space="preserve">{t}</m:t></m:r>'

def sub_(base, sub):
    return f'<m:sSub><m:e>{base}</m:e><m:sub>{sub}</m:sub></m:sSub>'

def sup_(base, exp):
    return f'<m:sSup><m:e>{base}</m:e><m:sup>{exp}</m:sup></m:sSup>'

def subsup_(base, sub, exp):
    return f'<m:sSubSup><m:e>{base}</m:e><m:sub>{sub}</m:sub><m:sup>{exp}</m:sup></m:sSubSup>'

def frac_(num, den):
    return f'<m:f><m:num>{num}</m:num><m:den>{den}</m:den></m:f>'

def sqrt_(content):
    return f'<m:rad><m:radPr><m:degHide m:val="1"/></m:radPr><m:deg/><m:e>{content}</m:e></m:rad>'

def paren_(content):
    return f'<m:d><m:dPr><m:begChr m:val="("/><m:endChr m:val=")"/><m:ctrlPr/></m:dPr><m:e>{content}</m:e></m:d>'

def bracket_(content):
    return f'<m:d><m:dPr><m:begChr m:val="["/><m:endChr m:val="]"/><m:ctrlPr/></m:dPr><m:e>{content}</m:e></m:d>'

def abs_(content):
    return f'<m:d><m:dPr><m:begChr m:val="|"/><m:endChr m:val="|"/><m:ctrlPr/></m:dPr><m:e>{content}</m:e></m:d>'

def norm_(content):
    return f'<m:d><m:dPr><m:begChr m:val="‖"/><m:endChr m:val="‖"/><m:ctrlPr/></m:dPr><m:e>{content}</m:e></m:d>'

def overbar_(content):
    return f'<m:acc><m:accPr><m:chr m:val="̅"/><m:ctrlPr/></m:accPr><m:e>{content}</m:e></m:acc>'

def func_(name, arg):
    return (f'<m:func><m:funcPr><m:ctrlPr/></m:funcPr>'
            f'<m:fName>{rp(name)}</m:fName><m:e>{arg}</m:e></m:func>')

def nary_(sym, sub_c, arg):
    """N-ary operator (sum/product) with subscript only."""
    return (f'<m:nary><m:naryPr><m:chr m:val="{sym}"/><m:limLoc m:val="subSup"/>'
            f'<m:subHide m:val="0"/><m:supHide m:val="1"/><m:ctrlPr/></m:naryPr>'
            f'<m:sub>{sub_c}</m:sub><m:sup/><m:e>{arg}</m:e></m:nary>')

def nary_sup_(sym, sub_c, sup_c, arg):
    """N-ary operator with subscript AND superscript."""
    return (f'<m:nary><m:naryPr><m:chr m:val="{sym}"/><m:limLoc m:val="subSup"/>'
            f'<m:ctrlPr/></m:naryPr>'
            f'<m:sub>{sub_c}</m:sub><m:sup>{sup_c}</m:sup><m:e>{arg}</m:e></m:nary>')

def norm_F_sq(content):
    """‖content‖²_F"""
    return subsup_(norm_(content), r('F'), r('2'))

# ═════════════════════════════════════════════════════════════════════════════
# Slide drawing helpers
# ═════════════════════════════════════════════════════════════════════════════

def add_text(slide, text, left, top, width, height,
             size=Pt(12), bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tx

def add_line(slide, y):
    shape = slide.shapes.add_shape(1, Inches(0.4), y, SW - Inches(0.8), Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.fill.background()

def add_eq(slide, omml_lines, left, top, width, height):
    """Embed one or more OMML equations into a textbox paragraph."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    txBody = tx.text_frame._txBody
    for p in txBody.findall(f'{{{NS_A}}}p'):
        txBody.remove(p)
    math_blocks = ''.join(f'<m:oMath>{ln}</m:oMath>' for ln in omml_lines)
    xml = (f'<a:p xmlns:a="{NS_A}" xmlns:a14="{NS_A14}" xmlns:m="{NS_M}">'
           f'<a14:m><m:oMathPara>{math_blocks}</m:oMathPara></a14:m></a:p>')
    txBody.append(etree.fromstring(xml))
    return tx

def slide_header(slide, title, subtitle=None):
    add_text(slide, title, Inches(0.4), Inches(0.18), SW - Inches(0.8), Inches(0.55),
             size=Pt(22), bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.73), SW - Inches(0.8), Inches(0.26),
                 size=Pt(10), color=GREY, italic=True)
    add_line(slide, Inches(0.98))

# ═════════════════════════════════════════════════════════════════════════════
# OMML equation definitions
# ═════════════════════════════════════════════════════════════════════════════

Fe   = sub_(r('F'), r('e'))
Je   = sub_(r('J'), r('e'))
Ce   = sub_(r('C'), r('e'))
Ae   = sub_(r('A'), r('e'))
Aebar= sub_(overbar_(r('A')), r('e'))
we   = sub_(r('w'), r('e'))
lk   = sub_(r('ℓ'), r('k'))   # ℓ_k
gamma_e = sub_(r('γ'), r('e'))
sigma_i = sub_(r('σ'), r('i'))

# F_e(t) = dx · DX^{-1}
eq_F = Fe + r('(t) = dx ⋅ ') + sup_(r('DX'), r('−1'))

# J_e(t) = det(F_e(t))
eq_J = Je + r('(t) = ') + func_('det', paren_(Fe + r('(t)')))

# Ā_e(t) = A_e(t) / A_e(0)
eq_Abar = Aebar + r('(t) = ') + frac_(Ae + r('(t)'), Ae + r('(0)'))

# C_e(t) = |J_e|^{-2/3} · F_e^T F_e
eq_C = (Ce + r('(t) = ') +
        sup_(abs_(Je), r('−') + frac_(r('2'), r('3'))) +
        r(' ⋅ ') + sup_(Fe, r('T')) + Fe)

# γ_e(t) = ½ tr(F_e^T F_e)
eq_gamma = (gamma_e + r('(t) = ') + frac_(r('1'), r('2')) + r(' ') +
            func_('tr', paren_(sup_(Fe, r('T')) + Fe)))

# GLE_e(t) = tr(½(F_e^T F_e − I)) = γ_e − 1
eq_GLE = (sub_(r('GLE'), r('e')) + r('(t) = ') +
          func_('tr', paren_(frac_(r('1'), r('2')) +
                paren_(sup_(Fe, r('T')) + Fe + r(' − I')))) +
          r(' = ') + gamma_e + r(' − 1'))

# EDI_e(t) = (1/3) Σ_{k=1}^{3} ℓ_k(t)/ℓ_k(0)
eq_EDI = (sub_(r('EDI'), r('e')) + r('(t) = ') + frac_(r('1'), r('3')) + r(' ') +
          nary_sup_('∑', r('k=1'), r('3'), frac_(lk + r('(t)'), lk + r('(0)'))))

# η(t) = 1 − std_e(Ā_e(t)) / σ_i ,  σ_i = ½√(N/(N−1))
eq_eta = (r('η(t) = 1 − ') +
          frac_(sub_(rp('std'), r('e')) + paren_(Aebar + r('(t)')), sigma_i) +
          r('   ,   where   ') +
          sigma_i + r(' = ') + frac_(r('1'), r('2')) + sqrt_(frac_(r('N'), r('N−1'))))

# ── Slide 2 ──

# Volumetric part of w: [(J^m-1)²/m² + (J^{-m}-1)²/m²]
vol = (frac_(sup_(paren_(sup_(Je, r('m')) + r(' − 1')), r('2')), sup_(r('m'), r('2'))) +
       r(' + ') +
       frac_(sup_(paren_(sup_(Je, r('−m')) + r(' − 1')), r('2')), sup_(r('m'), r('2'))))

# Isochoric part: [‖C^n−I‖²_F/n² + ‖C^{-n}−I‖²_F/n²]
iso = (frac_(norm_F_sq(sup_(Ce, r('n')) + r(' − I')), sup_(r('n'), r('2'))) +
       r(' + ') +
       frac_(norm_F_sq(sup_(Ce, r('−n')) + r(' − I')), sup_(r('n'), r('2'))))

eq_w1 = we + r('(t) = ') + frac_(r('K'), r('4')) + r(' ⋅ ') + bracket_(vol)
eq_w2 = r('          + ') + frac_(r('G'), r('4')) + r(' ⋅ ') + bracket_(iso)

eq_wmean = (sub_(r('w'), r('mean')) + r('(t) = ') + frac_(r('1'), r('N')) +
            nary_('∑', r('e'), we + r('(t)')))

eq_wstd = (sub_(r('w'), r('std')) + r('(t) = ') +
           sub_(rp('std'), r('e')) + paren_(we + r('(t)')))

eq_wCV = (sub_(r('w'), r('CV')) + r('(t) = ') +
          frac_(sub_(r('w'), r('std')) + r('(t)'),
                abs_(sub_(r('w'), r('mean')) + r('(t)')) + r(' + ε')))

eq_W = (r('W(t) = ') +
        nary_('∑', r('e'), we + r('(t) ⋅ ') + Ae + r('(0)')))

wbarA = sub_(overbar_(r('w')), r('A'))
Atot  = sub_(r('A'), r('tot'))

eq_wbarA = wbarA + r('(t) = ') + frac_(r('W(t)'), Atot)

eq_wstd_AW = (sub_(r('w'), r('std,AW')) + r('(t) = ') +
              sqrt_(frac_(
                  nary_('∑', r('e'),
                        Ae + r('(0) ⋅ ') +
                        sup_(paren_(we + r(' − ') + wbarA), r('2'))),
                  Atot)))

eq_wCV_AW = (sub_(r('w'), r('CV,AW')) + r('(t) = ') +
             frac_(sub_(r('w'), r('std,AW')) + r('(t)'),
                   abs_(wbarA + r('(t)')) + r(' + ε')))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  Geometric & Kinematic Parameters
# ═════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
slide_header(s1,
             "T2 Deformation Parameters — Geometry & Kinematics",
             "video_1002 & video_1003  ·  derived from triangle node coordinates only (no stress data)")

LX  = Inches(0.4)    # label x
EX  = Inches(2.95)   # equation x
LW  = Inches(2.45)   # label width
EW  = SW - EX - Inches(0.3)  # equation width
RH  = Inches(0.75)   # row height
Y0  = Inches(1.08)
LS  = Pt(11)

rows_s1 = [
    ("Deformation Gradient  F",    [eq_F]),
    ("Jacobian  J",                [eq_J]),
    ("Normalised Area  Ā",    [eq_Abar]),
    ("Isochoric Cauchy-Green  C",  [eq_C]),
    ("Shear  γ",              [eq_gamma]),
    ("Green-Lagrange  GLE",        [eq_GLE]),
    ("Edge Deform. Index  EDI",    [eq_EDI]),
    ("Mesh Regularity  η",    [eq_eta]),
]

for i, (label, eqs) in enumerate(rows_s1):
    y = Y0 + i * RH
    add_text(s1, label, LX, y + Inches(0.1), LW, RH, size=LS, bold=True)
    add_eq(s1, eqs, EX, y, EW, RH)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  Strain Energy Density & Statistics
# ═════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_header(s2,
             "Strain Energy Density & Statistics",
             "Parameters (K, G, m, n) = (1, 1, 2, 2)  ·  colourmap on mesh = wₑ per element  ·  time-series = global statistics")

# w_e formula (2 display lines)
add_text(s2, "wₑ — Strain Energy Density per element",
         LX, Inches(1.08), SW - Inches(0.8), Inches(0.28), size=LS, bold=True)
add_eq(s2, [eq_w1, eq_w2], LX, Inches(1.35), SW - Inches(0.8), Inches(1.55))

# Divider
add_text(s2, "Global statistics  (time-series vs. displacement)",
         LX, Inches(2.93), SW - Inches(0.8), Inches(0.28), size=LS, bold=True)
add_line(s2, Inches(3.19))

# 3-column grid for 6 statistics
CW  = (SW - Inches(0.8) - Inches(0.4)) / 3   # column width
CGX = [Inches(0.4), Inches(0.4) + CW + Inches(0.2), Inches(0.4) + 2*CW + Inches(0.4)]
LH  = Inches(0.26)
EH1 = Inches(0.75)
EH2 = Inches(1.3)   # taller row for 2-line stats
GY1 = Inches(3.27)
GY2 = GY1 + LH + EH1 + Inches(0.18)

stats = [
    # (col_x, y_label, label, eq_lines, eq_height)
    (CGX[0], GY1, "Mean energy  wₘₑₐₙ",           [eq_wmean],          EH1),
    (CGX[1], GY1, "Standard deviation  wₛₜₑ",          [eq_wstd],           EH1),
    (CGX[2], GY1, "Coeff. of variation  wᴄᴠ",              [eq_wCV],            EH1),
    (CGX[0], GY2, "Total area-weighted energy  W",                    [eq_W],              EH1),
    (CGX[1], GY2, "Area-weighted mean  w̅ₐ",               [eq_wbarA],          EH1),
    (CGX[2], GY2, "Area-weighted spread",                             [eq_wstd_AW, eq_wCV_AW], EH2),
]

for (cx, cy, label, eqs, eh) in stats:
    add_text(s2, label, cx, cy, CW, LH, size=Pt(10), bold=True)
    add_eq(s2, eqs, cx, cy + LH, CW, eh)

# ── Save ─────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(__file__), "T2_slides.pptx")
prs.save(OUT)
print(f"Saved → {OUT}")
